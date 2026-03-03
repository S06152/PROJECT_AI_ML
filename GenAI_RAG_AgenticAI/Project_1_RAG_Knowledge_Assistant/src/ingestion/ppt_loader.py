# Standard library imports
import sys
from src.utils.logger import logging
from src.utils.exception import CustomException
from typing import List
from langchain_core.documents import Document
from pptx import Presentation

class PPTLoader:
    """
    Wrapper class for loading PPTX (PowerPoint) files into LangChain Document objects.

    This class:
    - Accepts an uploaded PPTX file (e.g., from Streamlit).
    - Saves it temporarily to disk (since loader expects a file path).
    - Extracts slide content using UnstructuredPowerPointLoader.
    - Deletes the temporary file after processing.
    """

    def __init__(self, uploaded_file):
        """
        Initialize with uploaded file object.

        Args:
            uploaded_file: File-like object (e.g., Streamlit uploaded file).
        """
        try:
            logging.info("Initializing PPTLoader.")
            self.uploaded_file = uploaded_file

        except Exception as e:
            logging.exception("Error during PPTLoader initialization.")
            raise CustomException(e, sys)

    def load_documents(self) -> List[Document]:
        try:
            logging.info("Loading PPTX using python-pptx.")

            presentation = Presentation(self.uploaded_file)

            documents = []

            for i, slide in enumerate(presentation.slides):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                full_text = "\n".join(slide_text).strip()

                if full_text:
                    documents.append(
                        Document(
                            page_content=full_text,
                            metadata={"slide_number": i + 1}
                        )
                    )

            logging.info(f"Extracted {len(documents)} slides from PPTX.")
            return documents

        except Exception as e:
            logging.exception("Error while loading PPTX.")
            raise CustomException(e, sys)